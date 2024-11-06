from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PresentationCreator:
    def __init__(self, images_path='downloaded_images', output_file='presentation.pptx'):
        self.images_path = images_path
        self.output_file = output_file
        self.prs = Presentation()

    def create_pptx(self, slides_data):
        for idx, slide_data in enumerate(slides_data['slides']):
            slide_layout = self.prs.slide_layouts[0 if idx == 0 else 1]
            slide = self.prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']

            if idx > 0:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
                for point in slide_data['points']:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(0, 0, 0)

                image_path = os.path.join(self.images_path, f"{slide_data['image_keyword']['search_term']}.jpg")
                if os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, Inches(6.0), Inches(1.1), height=Inches(4.0))

        self.prs.save(self.output_file)
        return self.output_file
