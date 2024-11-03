from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
import os
from pathlib import Path
import PyPDF2
import time
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import requests
import json


GOOGLE_API_KEY = 'AIzaSyAT8sDfuOOwSbl_3_8KZoWdZQgyKI-XWJA'
genai.configure(api_key=GOOGLE_API_KEY)

app = FastAPI()

class PDFInput(BaseModel):
    pdf_path: str

def extract_text_from_pdf(pdf_file):
    with open(pdf_file, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''.join([page.extract_text() for page in reader.pages])
    return text

def create_structured_prompt(text):
    return f"""
    Lütfen aşağıdaki metinden profesyonel bir sunum oluştur. Her slayt için, slaytın içeriğini en iyi temsil edecek bir tablo dışı görsel arama terimi belirle.
    JSON formatında yanıt ver.

    Metin: {text}

    Yanıtı aşağıdaki JSON formatında ver:
    {{
        "slides": [
            {{
                "title": "Başlık",
                "points": [
                    "Madde 1",
                    "Madde 2",
                    "Madde 3"
                ],
                "image_keyword": {{
                    "search_term": "Arama terimi",
                    "description": "Bu terimin seçilme nedeni"
                }}
            }}
        ]
    }}

    Her slayt için image_keyword belirleme kuralları:
    1. Slaytın ana konusunu en iyi yansıtan bir terim seç
    2. Terimi Türkçe olarak belirt
    3. Tablo içeren görselleri dışla, tablo kullanma
    4. Görsel tipini belirt (diagram, illustration, photo, etc.)
    5. Bilimsel/akademik görseller için sonuna "scientific" veya "academic" ekle

    Sunum şunları içermeli:
    1. Kapak slaytı (başlık ve alt başlık)
    2. İçindekiler/Genel Bakış
    3-8. Ana içerik slaytları (her birinde 3-5 madde)
    9. Özet/Sonuç slaytı
    """

def download_image(keyword, download_path):
    if not os.path.exists(download_path):
        os.makedirs(download_path)

    options = webdriver.ChromeOptions()
    options.add_argument('--headless')
    options.add_argument('--ignore-certificate-errors')
    options.add_argument('--ignore-ssl-errors')
    driver = webdriver.Chrome(options=options)

    try:
        driver.get("https://www.google.com/imghp")
        search_box = driver.find_element(By.NAME, "q")
        search_box.send_keys(keyword)
        search_box.send_keys(Keys.RETURN)
        time.sleep(2)

        thumbnail = WebDriverWait(driver, 10).until(
            EC.presence_of_element_located((By.CSS_SELECTOR, "div.H8Rx8c"))
        )

        img_element = thumbnail.find_element(By.TAG_NAME, "img")
        driver.execute_script("arguments[0].click();", img_element)
        time.sleep(2)

        large_image = WebDriverWait(driver, 5).until(
            EC.presence_of_element_located((By.CSS_SELECTOR, 'img[jsname="kn3ccd"]'))
        )

        img_url = large_image.get_attribute('src')

        if img_url and img_url.startswith('http'):
            response = requests.get(img_url, timeout=10, verify=False)
            if response.status_code == 200:
                file_path = os.path.join(download_path, f"{keyword}.jpg")
                with open(file_path, 'wb') as f:
                    f.write(response.content)
                return file_path

    except Exception as e:
        print(f"Görsel indirme hatası: {str(e)}")
        try:
            thumbnails = driver.find_elements(By.CSS_SELECTOR, "img.rg_i")
            if thumbnails:
                img_url = thumbnails[0].get_attribute('src')
                if img_url and img_url.startswith('http'):
                    response = requests.get(img_url, timeout=10, verify=False)
                    if response.status_code == 200:
                        file_path = os.path.join(download_path, f"{keyword}.jpg")
                        with open(file_path, 'wb') as f:
                            f.write(response.content)
                        return file_path
        except Exception as backup_error:
            print(f"Alternatif görsel indirme hatası: {str(backup_error)}")
        return None
    finally:
        driver.quit()

def get_presentation_content(text):
    generation_config = {
        "temperature": 0.7,
        "top_p": 0.9,
        "top_k": 40,
        "max_output_tokens": 8192,
    }

    model = genai.GenerativeModel(
        model_name="gemini-1.5-pro",
        generation_config=generation_config
    )

    prompt = create_structured_prompt(text)
    response = model.generate_content(prompt)

    try:
        content = json.loads(response.text)
        for slide in content['slides']:
            if 'image_keyword' in slide:
                search_term = slide['image_keyword']['search_term']
                search_term += " high quality educational"
                slide['image_keyword']['search_term'] = search_term
        return content
    except json.JSONDecodeError:
        cleaned_response = response.text.strip()
        cleaned_response = cleaned_response[cleaned_response.find('{'):cleaned_response.rfind('}') + 1]
        return json.loads(cleaned_response)

from pptx.enum.text import PP_ALIGN

def create_pptx(slides_data, images_path, output_file):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(6)

    for idx, slide_data in enumerate(slides_data['slides']):
        if idx == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            title_frame = title_shape.text_frame
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 51, 102)

        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.top = Inches(0.3)
            title_shape.left = Inches(0.4)
            title_shape.width = Inches(9.2)
            title_shape.height = Inches(0.8)
            title_shape.text = slide_data['title']
            title_frame = title_shape.text_frame
            paragraph = title_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 51, 102)

            content = slide.placeholders[1]
            content.top = Inches(1.3)
            content.left = Inches(0.4)
            content.width = Inches(5.5)
            content.height = Inches(4.5)

            text_frame = content.text_frame
            text_frame.clear()
            text_frame.paragraphs[0].space_before = Pt(0)
            text_frame.paragraphs[0].space_after = Pt(12)

            for point in slide_data['points']:
                p = text_frame.add_paragraph()
                p.text = point
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.level = 0
                p.space_before = Pt(6)
                p.space_after = Pt(6)

            if 'image_keyword' in slide_data:
                search_term = slide_data['image_keyword']['search_term']
                print(f"\nSlayt {idx + 1} için arama terimi: {search_term}")
                print(f"Açıklama: {slide_data['image_keyword']['description']}")

                image_path = download_image(search_term, images_path)
                if image_path:
                    left = Inches(6.0)
                    top = Inches(1.1)
                    height = Inches(4.0)
                    try:
                        pic = slide.shapes.add_picture(image_path, left, top, height=height)
                        aspect_ratio = pic.width / pic.height
                        new_width = height * aspect_ratio
                        if new_width > Inches(3.5):
                            pic.width = Inches(3.5)
                            pic.height = Inches(3.5) / aspect_ratio
                    except Exception as e:
                        print(f"Görsel ekleme hatası: {str(e)}")

    prs.save(output_file)
    return output_file

def create_presentation(pdf_file):
    text = extract_text_from_pdf(pdf_file)
    slides_data = get_presentation_content(text)
    output_file = create_pptx(
        slides_data,
        images_path='/Users/muhammedtalhabicak/PycharmProjects/BTK24/downloaded_images',
        output_file=f"/Users/muhammedtalhabicak/PycharmProjects/BTK24/{Path(pdf_file).stem}_presentation.pptx"
    )
    return output_file

@app.post("/create-presentation")
async def generate_presentation(pdf_input: PDFInput):
    pdf_path = pdf_input.pdf_path
    if not os.path.isfile(pdf_path):
        raise HTTPException(status_code=404, detail="PDF dosyası bulunamadı")

    try:
        result_file = create_presentation(pdf_path)
        return {"message": "Sunum başarıyla oluşturuldu", "file_path": result_file}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Hata oluştu: {str(e)}")

import urllib3
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

if __name__ == "__main__":
    pdf_file = "/Users/muhammedtalhabicak/PycharmProjects/BTK24/SindirimSistemi.pdf"
    result = create_presentation(pdf_file)
    print(result)