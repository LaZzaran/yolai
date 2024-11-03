from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import os
from services.image_service import ImageService

class PresentationService:
    def __init__(self):
        self.prs = None

    def create_slide_layout(self, slide_data: dict, idx: int):
        """Slayt düzeni oluşturur"""
        slide_layout = self.prs.slide_layouts[1] if idx > 0 else self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        if idx > 0 and 'points' in slide_data:
            try:
                content = slide.placeholders[1]
                tf = content.text_frame
                tf.clear()

                for point in slide_data['points']:
                    p = tf.add_paragraph()
                    p.text = point
                    p.alignment = PP_ALIGN.LEFT
            except IndexError:
                print(f"Uyarı: Slayt {idx} için içerik yer tutucusu bulunamadı")

        return slide

    def create_pptx(self, slides_data: dict, output_dir: str) -> str:
        """Sunumu oluşturur"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(6)

        images_path = os.path.join(output_dir, "images")
        os.makedirs(images_path, exist_ok=True)

        for idx, slide_data in enumerate(slides_data['slides']):
            slide = self.create_slide_layout(slide_data, idx)

            if 'image_keyword' in slide_data:
                image_path = ImageService.download_image(
                    slide_data['image_keyword']['search_term'],
                    images_path
                )
                if image_path:
                    try:
                        left = Inches(6.0)
                        top = Inches(1.1)
                        width = Inches(3.5)
                        pic = slide.shapes.add_picture(image_path, left, top, width=width)
                        aspect_ratio = pic.height / pic.width
                        pic.height = int(width * aspect_ratio)
                    except Exception as e:
                        print(f"Görsel ekleme hatası: {str(e)}")

        output_file = os.path.join(output_dir, "presentation.pptx")
        self.prs.save(output_file)
        return output_file

    @staticmethod
    def add_video_to_slide(slide, video_path):
        """Add video to PowerPoint slide"""
        if video_path and os.path.exists(video_path):
            left = Inches(6.0)
            top = Inches(1.1)
            width = Inches(3.5)
            height = Inches(3.5)

            slide.shapes.add_movie(
                video_path,
                left, top, width, height,
                poster_frame_image=None,
                mime_type='video/mp4'
            )
