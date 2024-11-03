from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from stabilityai.services.model_runner import run_model
from pathlib import Path
from pptx.enum.text import PP_ALIGN

def create_pptx(slides_data, output_file):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(6)

    for idx, slide_data in enumerate(slides_data['slides']):
        if idx == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            title_frame = title_shape.text_frame
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 51, 102)

        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.top = Inches(0.3)
            title_shape.left = Inches(0.4)
            title_shape.width = Inches(9.2)
            title_shape.height = Inches(0.8)
            title_shape.text = slide_data['title']
            title_frame = title_shape.text_frame
            paragraph = title_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 51, 102)

            content = slide.placeholders[1]
            content.top = Inches(1.3)
            content.left = Inches(0.4)
            content.width = Inches(5.5)
            content.height = Inches(4.5)

            text_frame = content.text_frame
            text_frame.clear()
            text_frame.paragraphs[0].space_before = Pt(0)
            text_frame.paragraphs[0].space_after = Pt(12)

            for point in slide_data['points']:
                p = text_frame.add_paragraph()
                p.text = point
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.level = 0
                p.space_before = Pt(6)
                p.space_after = Pt(6)

            if 'image_keyword' in slide_data:
                search_term = slide_data['image_keyword']['search_term']
                output_image_path = f"{Path(output_file).stem}_{search_term.replace(' ', '_')}.png"
                run_model(prompt=search_term, output_path=output_image_path) # Şu kısımda stabilityai dan image oluşturumu yapılıyor

                left = Inches(6.0)
                top = Inches(1.1)
                height = Inches(4.0)
                try:
                    pic = slide.shapes.add_picture(output_image_path, left, top, height=height)
                    aspect_ratio = pic.width / pic.height
                    new_width = height * aspect_ratio
                    if new_width > Inches(3.5):
                        pic.width = Inches(3.5)
                        pic.height = Inches(3.5) / aspect_ratio
                except Exception as e:
                    print(f"Görsel ekleme hatası: {str(e)}")

    prs.save(output_file)
    return output_file
